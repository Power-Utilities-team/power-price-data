"""
check_consistency.py — fail-loud guard that the LINKED deck and the STATIC deck
never drift from deck_spec.py (or each other). Run before every delivery.

Asserts, for BOTH decks:
  * slide count == 1 (title) + len(deck_spec.SLIDES)
  * each content slide's title + kicker == deck_spec
  * each content slide's navy-bar captions (in order) == deck_spec exhibit captions
And that the linked workbook holds charts 1..19.

Exit 0 = consistent; exit 1 = drift (prints the diffs).
"""
from __future__ import annotations
import os, sys, zipfile, re
from xml.etree import ElementTree as ET
from pptx import Presentation
import deck_spec

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
LINKED = os.path.join(ROOT, "outputs", "HourlyPowerData.pptx")
STATIC = os.path.join(ROOT, "outputs", "HourlyPowerData_snapshot.pptx")
WORKBOOK = os.path.join(ROOT, "outputs", "HourlyPowerData.xlsx")


n_charts_msg = "?"   # filled by main() once the workbook's expected count is known


def deck_content(path):
    """Return [(title, kicker, [captions...]), ...] for content slides (skip title)."""
    prs = Presentation(path)
    out = []
    for s in list(prs.slides)[1:]:
        title = kicker = ""
        caps = []
        for ph in s.placeholders:
            idx = ph.placeholder_format.idx
            if idx == 0: title = ph.text.strip()
            elif idx == 13: kicker = ph.text.strip()
        for sh in s.shapes:
            if not sh.is_placeholder and sh.has_text_frame and sh.text_frame.text.strip():
                caps.append(sh.text_frame.text.strip())
        out.append((title, kicker, caps))
    return out


def expected():
    return [(s["title"], s["kicker"], [e["caption"] for e in s["exhibits"]]) for s in deck_spec.SLIDES]


def check_deck(name, path, exp):
    errs = []
    if not os.path.exists(path):
        return [f"{name}: file missing ({path})"]
    got = deck_content(path)
    if len(got) != len(exp):
        errs.append(f"{name}: {len(got)} content slides, expected {len(exp)}")
    for i, (e, g) in enumerate(zip(exp, got), start=1):
        if e[0] != g[0]:
            errs.append(f"{name} slide {i} title: got {g[0]!r} != spec {e[0]!r}")
        if e[1] != g[1]:
            errs.append(f"{name} slide {i} kicker: got {g[1]!r} != spec {e[1]!r}")
        if e[2] != g[2]:
            errs.append(f"{name} slide {i} captions: got {g[2]} != spec {e[2]}")
    return errs


def check_xml_wellformed(path, label):
    """Every XML part must parse. This is the cheapest possible check and the one that
    matters most: a single unescaped '&' in one cell makes Excel offer to Recover the
    workbook, and Recovery strips Power Query. On 2026-07-22 a text column written into a
    numeric <v> shipped three malformed sheets through a PASSing consistency run, because
    nothing here had ever actually parsed the file it was signing off.
    """
    errs = []
    z = zipfile.ZipFile(path)
    for n in z.namelist():
        if not n.endswith((".xml", ".rels")):
            continue
        try:
            ET.fromstring(z.read(n))
        except ET.ParseError as e:
            errs.append(f"{label} {n}: malformed XML — {e}")
    return errs


def check_refresh_stability(path):
    """The two halves of the invariant that keeps charts intact through a refresh.

    Excel re-anchors a chart series whose range runs to or past the end of the data it
    reads. So (a) every query tab's pre-filled cache must be exactly the size of the CSV
    it will load, and (b) no chart range may extend past that cache. Break either and a
    chart silently re-fits itself on the user's first refresh — which is how chart12
    quietly re-acquired the six technologies curate_tech_charts.py had removed.
    """
    errs = []
    z = zipfile.ZipFile(path)
    parts = {n: z.read(n) for n in z.namelist()}
    wb = parts["xl/workbook.xml"].decode()
    relmap = dict(re.findall(r'Id="([^"]+)"[^>]*Target="([^"]+)"',
                             parts["xl/_rels/workbook.xml.rels"].decode()))
    extent = {}
    for name, rid in re.findall(r'<sheet name="([^"]+)"[^>]*?r:id="(rId\d+)"', wb):
        spart = "xl/" + relmap[rid].lstrip("/")
        sx = parts[spart].decode(errors="replace")
        d = re.search(r'<dimension ref="[A-Z]+\d+:[A-Z]+(\d+)"', sx)
        if d:
            extent[name] = int(d.group(1))
        rels = parts.get(spart.replace("worksheets/", "worksheets/_rels/") + ".rels", b"").decode()
        t = re.search(r'Target="\.\./tables/(table\d+\.xml)"', rels)
        if not t or name not in extent:
            continue
        tbl = parts["xl/tables/" + t.group(1)].decode()
        stem = re.search(r'\sname="([^"]+)"', tbl).group(1)
        csv_path = os.path.join(ROOT, "published", "charts", stem + ".csv")
        if not os.path.exists(csv_path):
            continue
        with open(csv_path, newline="") as f:
            want = sum(1 for _ in f)
        if extent[name] != want:
            errs.append(f"WORKBOOK {name}: pre-filled {extent[name]} rows but "
                        f"{stem}.csv has {want} — the table will change shape on refresh "
                        f"and Excel may re-anchor a chart (run resync_prefill.py)")

    for n in sorted(x for x in parts if re.match(r"xl/charts/chart\d+\.xml$", x)):
        cx = parts[n].decode()
        for f in sorted(set(re.findall(r"<c:f>([^<]+)</c:f>", cx))):
            m = re.match(r"([^!]+)!\$[A-Z]+\$\d+:\$[A-Z]+\$(\d+)$", f)
            if m and m.group(1) in extent and int(m.group(2)) > extent[m.group(1)]:
                errs.append(f"WORKBOOK {os.path.basename(n)}: range {f} runs past the "
                            f"pre-filled data (row {extent[m.group(1)]}) — Excel will "
                            f"stretch this series on refresh")
    return errs


def check_content_types(path, label):
    """Every part in the package must be declared in [Content_Types].xml.

    A part that is perfectly well-formed XML but undeclared makes the PACKAGE invalid.
    Excel then opens with "We found a problem with some content … do you want us to try
    to recover", and recovering strips Power Query — the one outcome this project treats
    as unacceptable.

    This is not hypothetical: on 2026-07-31 two new worksheets were added while the
    content-type list was still a hardcoded literal, so their overrides were omitted.
    Every existing check passed — the XML was well-formed, chart structure matched the
    spec — because they all inspect parts individually and none inspected the manifest.
    Only opening the file in Excel revealed it. Hence this gate.
    """
    errs = []
    z = zipfile.ZipFile(path)
    ct = z.read("[Content_Types].xml").decode()
    overrides = set(re.findall(r'PartName="([^"]+)"', ct))
    defaults = set(re.findall(r'Extension="([^"]+)"', ct))
    for n in z.namelist():
        if n == "[Content_Types].xml" or n.endswith("/"):
            continue
        ext = n.rsplit(".", 1)[-1].lower() if "." in n else ""
        if ("/" + n) in overrides:
            continue
        # These part families are TYPED: each needs its own Override. The package also
        # carries `Default xml -> application/xml`, so an undeclared worksheet does not
        # look missing — it silently inherits the WRONG type, which is exactly why the
        # 2026-07-31 build looked fine to every structural check and still made Excel
        # offer to Recover. A Default must therefore never satisfy these.
        if n.startswith(("xl/worksheets/", "xl/charts/", "xl/drawings/", "xl/tables/")) \
                and ext == "xml":
            errs.append(f"{label}: {n} has no Content_Types Override — it falls back to "
                        f"the generic xml Default, and Excel will offer to Recover "
                        f"(which strips Power Query)")
            continue
        if ext in defaults:
            continue        # genuinely covered by a Default — .rels, images, .bin
        errs.append(f"{label}: {n} has no Content_Types entry at all")
    return errs


def check_no_future_data():
    """No published feed may carry data for a period that has not finished.

    The project's own rule is that a period-based chart never shows a partial period,
    and every producer applies it — but only to its OWN output, so a gap in one
    producer was invisible to every check. G1's quarterly average was broadcast across
    the whole quarter and published 2.5 months into the future, while the PNG path
    clipped the same series correctly: the workbook and the deck drew one exhibit two
    ways. Structural checks could not see it because the chart XML was perfectly valid.

    This asserts the rule on the artefacts Excel actually loads, so the two update
    paths cannot silently diverge on data cutoffs again.
    """
    import csv
    from datetime import date
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
    from completeness import cutoffs

    c = cutoffs()
    qend = str(c["last_complete_quarter_end"])[:10]
    today = date.today().isoformat()
    pub = os.path.join(ROOT, "published", "charts")
    errs = []
    if not os.path.isdir(pub):
        return errs

    for fn in sorted(os.listdir(pub)):
        if not fn.endswith(".csv"):
            continue
        path = os.path.join(pub, fn)
        with open(path) as fh:
            rows = list(csv.reader(fh))
        if len(rows) < 2:
            continue
        header = rows[0]
        if not header or header[0].strip().lower() not in ("date", "day"):
            continue
        # quarterly-average columns are the ones gated to a completed quarter;
        # everything else on a daily axis must simply not run past today.
        qcols = [i for i, h in enumerate(header) if h.endswith("_qavg")]
        for r in rows[1:]:
            if not r or not r[0].strip():
                continue
            d = r[0].strip()[:10]
            if not re.match(r"\d{4}-\d{2}-\d{2}$", d):
                continue
            for i, cell in enumerate(r):
                if i == 0 or not cell.strip():
                    continue
                limit = qend if i in qcols else today
                if d > limit:
                    errs.append(f"{fn}: {header[i] if i < len(header) else i} has data "
                                f"at {d}, past {limit}")
                    break
            if len(errs) > 6:
                return errs + ["(further future-data errors suppressed)"]
    return errs


def main():
    exp = expected()
    errs = []
    errs += check_no_future_data()
    errs += check_deck("LINKED", LINKED, exp)
    errs += check_deck("STATIC", STATIC, exp)
    # The workbook holds a contiguous run of chart parts: the 19 the original chain
    # builds, then the per-country variants and the monthly capture exhibits that
    # add_extra_charts.py appends. The expected total is DERIVED from the same lists
    # those are built from, so adding a country or a charted technology moves the
    # assertion with the build instead of failing it. A hardcoded 19 is what this line
    # used to say, and it is why the check could not see that the circulated workbook
    # had 62 charts while the pipeline produced 19.
    if os.path.exists(WORKBOOK):
        global n_charts_msg
        import add_extra_charts as _extra
        # NOT named `expected`: that is a module-level function used a few lines above,
        # and shadowing it makes the earlier call fail with an UnboundLocalError that
        # points at the wrong line entirely.
        n_charts = n_charts_msg = _extra.expected_chart_total()
        z = zipfile.ZipFile(WORKBOOK)
        nums = sorted(int(re.search(r"chart(\d+)", n).group(1))
                      for n in z.namelist() if re.match(r"xl/charts/chart\d+\.xml$", n))
        if nums != list(range(1, n_charts + 1)):
            errs.append(f"WORKBOOK charts: {len(nums)} parts, expected {n_charts} "
                        f"(1..{n_charts}); got {nums[:3]}..{nums[-3:]}")
        errs += check_xml_wellformed(WORKBOOK, 'WORKBOOK')
        errs += check_content_types(WORKBOOK, 'WORKBOOK')
        errs += check_refresh_stability(WORKBOOK)
    else:
        errs.append(f"WORKBOOK missing ({WORKBOOK})")
    frozen = os.path.join(ROOT, 'outputs', 'HourlyPowerData_frozen.xlsx')
    if os.path.exists(frozen):
        errs += check_xml_wellformed(frozen, 'FROZEN')
        errs += check_content_types(frozen, 'FROZEN')

    if errs:
        print("CONSISTENCY: FAIL")
        for e in errs: print("  ✗", e)
        sys.exit(1)
    print(f"CONSISTENCY: PASS — both decks match deck_spec ({len(exp)} content slides, "
          f"{sum(len(s['exhibits']) for s in deck_spec.SLIDES)} exhibits) + workbook charts 1-{n_charts_msg}.")


if __name__ == "__main__":
    main()
